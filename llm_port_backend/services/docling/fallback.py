"""Lightweight fallback document text extraction.

Used when the Docling microservice is **not enabled** or unreachable.
Provides basic text extraction using small, dependency-light libraries:

* **PDF** → ``pdfplumber`` — text extraction with basic table
  detection, no OCR.
* **DOCX** → ``python-docx`` — paragraph / table text.
* **PPTX** → ``python-pptx`` — slide text frames.
* **XLSX / CSV** → ``openpyxl`` / stdlib ``csv`` — cell values.
* **HTML** → stdlib ``html.parser`` — stripped text.
* **Plain text** (TXT, MD, AsciiDoc, XML) — read as-is.

The quality is significantly lower than Docling (no OCR, no table
structure, no heading hierarchy), but it works without a GPU and
without an external service.
"""

from __future__ import annotations

import csv
import io
import logging
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ── Result dataclass ──────────────────────────────────────────────────


class FallbackResult:
    """Mirror of the Docling ConvertResponse shape."""

    __slots__ = ("content", "metadata", "chunks")

    def __init__(
        self,
        content: str,
        *,
        page_count: int = 1,
        format_detected: str = "unknown",
        tables_found: int = 0,
    ) -> None:
        self.content = content
        self.metadata = {
            "page_count": page_count,
            "format_detected": format_detected,
            "tables_found": tables_found,
            "ocr_applied": False,
            "processing_time_ms": 0,
        }
        self.chunks: list[dict[str, Any]] = []

    def to_dict(self) -> dict[str, Any]:
        return {
            "content": self.content,
            "metadata": self.metadata,
            "chunks": self.chunks,
        }


# ── Format-specific extractors ───────────────────────────────────────


def _extract_pdf(data: bytes) -> FallbackResult:
    """Extract text from PDF using pdfplumber."""
    try:
        import pdfplumber  # noqa: PLC0415
    except ImportError:
        logger.warning("pdfplumber is not installed — cannot extract PDF text")
        return FallbackResult("", format_detected="pdf")

    pages: list[str] = []
    tables_found = 0
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            parts: list[str] = []

            # Extract tables first so we can include structured data
            page_tables = page.extract_tables() or []
            tables_found += len(page_tables)
            table_regions: list[Any] = []
            for tbl in page_tables:
                rows = [
                    " | ".join(str(c) if c else "" for c in row)
                    for row in tbl
                    if any(row)
                ]
                if rows:
                    parts.append("\n".join(rows))
                # Collect bounding boxes so we can exclude table text
                # from the main text extraction (avoid duplication)
                for t in page.find_tables():
                    table_regions.append(t.bbox)

            # Extract remaining (non-table) text
            if table_regions:
                # Crop out table areas to avoid duplicate text
                cropped = page
                for bbox in table_regions:
                    cropped = cropped.outside_bbox(bbox)
                text = (cropped.extract_text() or "").strip()
            else:
                text = (page.extract_text() or "").strip()

            if text:
                parts.insert(0, text)

            if parts:
                pages.append("\n\n".join(parts))

    return FallbackResult(
        "\n\n".join(pages),
        page_count=len(pages),
        format_detected="pdf",
        tables_found=tables_found,
    )


def _extract_docx(data: bytes) -> FallbackResult:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document  # noqa: PLC0415
    except ImportError:
        logger.warning("python-docx is not installed — cannot extract DOCX text")
        return FallbackResult("", format_detected="docx")

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    tables_found = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        tables_found += 1
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        parts.append("\n".join(rows))

    return FallbackResult(
        "\n\n".join(parts),
        format_detected="docx",
        tables_found=tables_found,
    )


def _extract_pptx(data: bytes) -> FallbackResult:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        logger.warning("python-pptx is not installed — cannot extract PPTX text")
        return FallbackResult("", format_detected="pptx")

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []

    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        if parts:
            slides.append("\n".join(parts))

    return FallbackResult(
        "\n\n---\n\n".join(slides),
        page_count=len(slides),
        format_detected="pptx",
    )


def _extract_xlsx(data: bytes) -> FallbackResult:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    except ImportError:
        logger.warning("openpyxl is not installed — cannot extract XLSX text")
        return FallbackResult("", format_detected="xlsx")

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets: list[str] = []
    tables_found = 0

    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            tables_found += 1
            sheets.append(f"## {ws.title}\n\n" + "\n".join(rows))
    wb.close()

    return FallbackResult(
        "\n\n".join(sheets),
        page_count=len(sheets),
        format_detected="xlsx",
        tables_found=tables_found,
    )


def _extract_csv(data: bytes) -> FallbackResult:
    """Extract text from CSV using stdlib csv."""
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(row) for row in reader if any(row)]
    return FallbackResult(
        "\n".join(rows),
        format_detected="csv",
        tables_found=1 if rows else 0,
    )


class _HTMLTextExtractor(HTMLParser):
    """Simple HTML → plain text via stdlib."""

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        self._parts.append(data)

    @property
    def text(self) -> str:
        return "".join(self._parts).strip()


def _extract_html(data: bytes) -> FallbackResult:
    """Strip HTML tags and return plain text."""
    parser = _HTMLTextExtractor()
    parser.feed(data.decode("utf-8", errors="replace"))
    return FallbackResult(parser.text, format_detected="html")


def _extract_plaintext(data: bytes, fmt: str = "txt") -> FallbackResult:
    """Return file contents as-is (TXT, MD, AsciiDoc, XML, etc.)."""
    return FallbackResult(
        data.decode("utf-8", errors="replace"),
        format_detected=fmt,
    )


# ── Dispatcher ────────────────────────────────────────────────────────

_EXT_MAP: dict[str, Any] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".csv": _extract_csv,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".txt": lambda d: _extract_plaintext(d, "txt"),
    ".md": lambda d: _extract_plaintext(d, "md"),
    ".adoc": lambda d: _extract_plaintext(d, "adoc"),
    ".xml": lambda d: _extract_plaintext(d, "xml"),
    ".json": lambda d: _extract_plaintext(d, "json"),
    ".rst": lambda d: _extract_plaintext(d, "rst"),
}


def extract_text(
    file_bytes: bytes,
    filename: str,
) -> FallbackResult:
    """Extract text from ``file_bytes`` using lightweight libraries.

    Falls back to plain-text decoding for unknown extensions.
    """
    ext = Path(filename).suffix.lower()
    handler = _EXT_MAP.get(ext, lambda d: _extract_plaintext(d, ext.lstrip(".")))
    try:
        return handler(file_bytes)
    except Exception:
        logger.exception("Fallback extraction failed for %s", filename)
        return FallbackResult(
            "",
            format_detected=ext.lstrip(".") or "unknown",
        )
